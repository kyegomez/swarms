from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define slide layout (1: Title Slide, 5: Title and Content)
slide_layout_title = prs.slide_layouts[0]
slide_layout_content = prs.slide_layouts[1]


# Function to add title slide
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(slide_layout_title)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

    # Set title style
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = (
        RGBColor(255, 0, 0)
    )  # Red text

    # Set subtitle style
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = (
        RGBColor(255, 255, 255)
    )  # White text


# Function to add content slides
def add_content_slide(title, content):
    slide = prs.slides.add_slide(slide_layout_content)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

    # Set title style
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = (
        RGBColor(255, 0, 0)
    )  # Red text

    # Set content style
    content_text_frame = content_placeholder.text_frame
    content_text_frame.word_wrap = True

    for paragraph in content_text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(
            255, 255, 255
        )  # White text
        paragraph.alignment = PP_ALIGN.LEFT


# Adding slides according to your agenda

# Title slide
add_title_slide(
    "Harnessing Thousands of Agents",
    "Automating Accounting, Marketing, and Beyond with the Spreadsheet Swarm",
)

# Introduction Slide
add_content_slide(
    "Introduction to Spreadsheet Swarm",
    "Overview of Swarms and their application in automating business operations.\n\nTime: 15 mins",
)

# Automating Accounting Slide
add_content_slide(
    "Automating Accounting",
    "How to mass analyze client transactions and data analysis using the Spreadsheet Swarm.\n\nTime: 25 mins",
)

# Automating Marketing Operations Slide
add_content_slide(
    "Automating Marketing Operations",
    "Scaling marketing campaigns, customer segmentation, and content generation with Swarms.\n\nTime: 25 mins",
)

# Automating Finance Operations Slide
add_content_slide(
    "Automating Finance Operations",
    "Automating financial forecasting, transaction analysis, and report generation with ease.\n\nTime: 25 mins",
)

# Live Demo and Q&A Slide
add_content_slide(
    "Live Demo and Q&A",
    "Walkthrough of real-world examples and a live demonstration of the Spreadsheet Swarm in action.\n\nOpen floor for questions and discussions.\n\nTime: 30 mins",
)

# Save the presentation
prs.save("Harnessing_Thousands_of_Agents.pptx")
